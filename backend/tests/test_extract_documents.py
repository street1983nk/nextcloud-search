"""PDF, OOXML and OpenDocument, including every way each of them can go wrong.

The reference corpus carries the four PDF cases that matter, so they are read
from there rather than rebuilt: those files are what the read only gate copies
into a throwaway Nextcloud, and a test that invents its own PDF would stop
saying anything about the file the gate actually sees.

Everything the corpus does not contain is built here, in the test, from the
standard library. A multi page PDF, a spreadsheet over the cell limit and three
OpenDocument files are all about one specific shape, and adding them to the
corpus would grow a fixture set that exists for a different purpose. Built in
the test they are visible in the diff, which is where a reviewer looks.
"""

from __future__ import annotations

import time
from collections.abc import Callable, Sequence
from pathlib import Path
from zipfile import ZipFile

import docx
import openpyxl
import pypdfium2
import pytest
from lxml import etree  # pyright: ignore[reportAttributeAccessIssue]
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Emu
from pypdf.errors import EmptyFileError, FileNotDecryptedError

from findling import config
from findling.extract.dispatch import ALLOWED_MIMETYPES, Route, extract
from findling.extract.errors import ExtractionOutcome, Reason, State
from findling.extract.odf import extract_odf
from findling.extract.office import extract_docx, extract_pptx, extract_xlsx
from findling.extract.pdf import _MIN_CHARS_PER_PAGE, _SCAN_PAGE_SHARE, extract_pdf

CORPUS = Path(__file__).resolve().parents[2] / "testdata" / "corpus"


def _write(directory: Path, name: str, payload: bytes) -> str:
    target = directory / name
    target.write_bytes(payload)
    return str(target)


def _assemble_pdf(objects: list[bytes]) -> bytes:
    """Numbered objects plus a correct cross reference table, as in build_corpus.py.

    Deliberately the same shape as the corpus builder: a PDF whose xref table is
    wrong is a corrupt PDF, and a test fixture that is accidentally corrupt would
    prove the error path while claiming to prove the happy one.
    """
    out = bytearray(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode("ascii") + body + b"\nendobj\n"

    xref_offset = len(out)
    size = len(objects) + 1
    out += f"xref\n0 {size}\n".encode("ascii")
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode("ascii")
    out += f"trailer\n<< /Size {size} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    return bytes(out)


def _multi_page_pdf(pages: int) -> bytes:
    """A PDF of the given length, every page carrying a full line of real text.

    The lines are long on purpose. A page with three words would be judged as a
    page without a text layer, and this fixture is about the page cap, not about
    the text layer threshold.
    """
    font_number = 3 + 2 * pages
    kids = " ".join(f"{3 + 2 * index} 0 R" for index in range(pages))
    objects: list[bytes] = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        f"<< /Type /Pages /Kids [{kids}] /Count {pages} >>".encode("ascii"),
    ]
    for index in range(pages):
        content = (
            f"BT /F1 10 Tf 20 70 Td (Sitzungsvorlage der Gemeinde, Blatt {index + 1}) Tj ET\n"
            f"BT /F1 10 Tf 20 45 Td (Betreff: Grundstuecksverkehrsgenehmigung) Tj ET\n"
        ).encode("ascii")
        objects.append(
            (
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 400 120]"
                f" /Resources << /Font << /F1 {font_number} 0 R >> >> /Contents {4 + 2 * index} 0 R >>"
            ).encode("ascii")
        )
        objects.append(f"<< /Length {len(content)} >>\nstream\n".encode("ascii") + content + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    return _assemble_pdf(objects)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def test_a_pdf_with_a_text_layer_yields_its_embedded_text() -> None:
    outcome = extract_pdf(str(CORPUS / "01-text-layer.pdf"))

    assert outcome.state is State.INDEXED
    assert "Findling reference corpus" in outcome.text
    assert "real text layer" in outcome.text


def test_a_pdf_without_a_text_layer_becomes_the_ocr_queue_of_phase_three() -> None:
    # Not a failure and not empty_text: this exact verdict is the list phase 3
    # works through, and a wrong bucket here costs a full reindex later.
    outcome = extract_pdf(str(CORPUS / "02-scan-no-text-layer.pdf"))

    assert outcome == ExtractionOutcome.skipped(Reason.NO_TEXT_LAYER)


def test_a_password_protected_pdf_is_answered_before_pdfium_is_asked(monkeypatch: pytest.MonkeyPatch) -> None:
    # The order is the whole point. pypdf answers the encryption question without
    # touching the pages, and reading .pages on a protected file raises, which
    # would turn a deliberate decision into a failure. pdfium is replaced by
    # something that cannot be called, so the claim is proven and not assumed.
    def unreachable(*_args: object, **_kwargs: object) -> object:
        raise AssertionError("pdfium was opened although the file is encrypted")

    monkeypatch.setattr(pypdfium2, "PdfDocument", unreachable)

    outcome = extract_pdf(str(CORPUS / "07-password-protected.pdf"))

    assert outcome == ExtractionOutcome.skipped(Reason.ENCRYPTED)


def test_a_zero_byte_pdf_is_failed_empty_file() -> None:
    outcome = extract_pdf(str(CORPUS / "06-zero-bytes.pdf"))

    assert outcome == ExtractionOutcome.failed(Reason.EMPTY_FILE)


def test_a_pdf_header_followed_by_garbage_is_failed_corrupt(tmp_path: Path) -> None:
    path = _write(tmp_path, "kaputt.pdf", b"%PDF-1.7\n" + b"garbage" * 50)

    outcome = extract_pdf(path)

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


def test_page_and_textpage_are_closed_even_when_a_page_raises(monkeypatch: pytest.MonkeyPatch) -> None:
    # pdfium hands out C resources. Over a run of 100000 files an object that is
    # never closed is a leak, and the error path is where closing gets forgotten.
    closed: list[str] = []
    close_page = pypdfium2.PdfPage.close
    close_textpage = pypdfium2.PdfTextPage.close

    def note_page(self: pypdfium2.PdfPage) -> None:
        closed.append("page")
        close_page(self)

    def note_textpage(self: pypdfium2.PdfTextPage) -> None:
        closed.append("textpage")
        close_textpage(self)

    def explode(self: pypdfium2.PdfTextPage) -> str:
        raise RuntimeError("pdfium lost its footing halfway through a page")

    monkeypatch.setattr(pypdfium2.PdfPage, "close", note_page)
    monkeypatch.setattr(pypdfium2.PdfTextPage, "close", note_textpage)
    monkeypatch.setattr(pypdfium2.PdfTextPage, "get_text_bounded", explode)

    with pytest.raises(RuntimeError):
        extract_pdf(str(CORPUS / "01-text-layer.pdf"))

    assert closed == ["textpage", "page"]


def test_a_pdf_over_the_page_cap_stops_and_says_truncated(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    monkeypatch.setenv("FINDLING_MAX_PDF_PAGES", "2")
    config.settings.cache_clear()
    try:
        path = _write(tmp_path, "lang.pdf", _multi_page_pdf(4))

        outcome = extract_pdf(path)

        assert outcome.state is State.INDEXED
        assert outcome.truncated is True
        assert "Blatt 2" in outcome.text
        assert "Blatt 3" not in outcome.text
    finally:
        config.settings.cache_clear()


def test_a_pypdf_empty_file_error_that_escapes_keeps_its_own_reason() -> None:
    # The extractor catches this one itself. The table entry is the net for the
    # day somebody reorders that function: an escaped zero byte error stays
    # empty_file instead of falling into the blanket corrupt.
    outcome = ExtractionOutcome.from_exception(EmptyFileError("nothing to read"))

    assert outcome == ExtractionOutcome.failed(Reason.EMPTY_FILE)


def test_an_encrypted_pdf_is_never_translated_out_of_an_exception() -> None:
    # failed(encrypted) is not a pair the taxonomy has, so a table entry for
    # FileNotDecryptedError would raise inside the error handler instead of
    # producing a verdict. The decision belongs where the file is opened.
    outcome = ExtractionOutcome.from_exception(FileNotDecryptedError("password required"))

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


# ---------------------------------------------------------------------------
# The text layer threshold, measured over the corpus on 2026-09-01.
#
# The numbers these tests lean on are in docs/ocr.md under "Die Textlayer-
# Erkennung": a full A4 page of German administrative prose measures 442 to 456
# characters, the sparsest genuine text page of the corpus measures 29, a page
# whose only text object is a headline measures 12, and every one of the nine
# rendered scan pages measures exactly 0.
# ---------------------------------------------------------------------------

_LINE_WIDTH = 38


def _filler(count: int) -> str:
    """A line of exactly ``count`` characters, which is what pdfium hands back."""
    base = "Sitzungsvorlage der Gemeinde Musterhausen "
    text = (base * (count // len(base) + 1))[:count]
    return text[:-1] + "x" if text.endswith(" ") else text


def _lines_for(characters: int) -> tuple[str, ...]:
    """``characters`` spread over lines that fit on the page.

    One long line would run off the media box, and get_text_bounded only reports
    what stands inside it: the fixture would then measure the page width instead
    of the threshold.
    """
    full, rest = divmod(characters, _LINE_WIDTH)
    lines = [_filler(_LINE_WIDTH) for _ in range(full)]
    if rest:
        lines.append(_filler(rest))
    return tuple(lines)


def _pdf_with_page_lines(pages: Sequence[Sequence[str]]) -> bytes:
    """A PDF where every page carries exactly the given lines, empty pages included.

    An empty page here is the shape of a scanned page for this question: a page
    object with a content stream that draws no text. What a real scan puts on it
    instead is an image, and the corpus carries those; a fixture that has to vary
    the character count per page does not need the pixels.
    """
    font_number = 3 + 2 * len(pages)
    kids = " ".join(f"{3 + 2 * index} 0 R" for index in range(len(pages)))
    objects: list[bytes] = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode("ascii"),
    ]
    for index, lines in enumerate(pages):
        content = bytearray()
        baseline = 650
        for line in lines:
            content += f"BT /F1 10 Tf 20 {baseline} Td ({line}) Tj ET\n".encode("ascii")
            baseline -= 25
        objects.append(
            (
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 400 700]"
                f" /Resources << /Font << /F1 {font_number} 0 R >> >> /Contents {4 + 2 * index} 0 R >>"
            ).encode("ascii")
        )
        objects.append(f"<< /Length {len(content)} >>\nstream\n".encode("ascii") + bytes(content) + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    return _assemble_pdf(objects)


def test_a_scanned_council_document_is_handed_over_as_no_text_layer() -> None:
    # Three A4 pages of German prose that exist only as pixels. Every one of them
    # measures zero characters, so the whole file is the OCR queue and nothing
    # else. This is the file the acceptance of D-09 later searches for.
    outcome = extract_pdf(str(CORPUS / "13-ratsvorlage-scan.pdf"))

    assert outcome == ExtractionOutcome.skipped(Reason.NO_TEXT_LAYER)


def test_mixed_pdf_with_text_layer_and_scanned_annex_is_extracted() -> None:
    # Pitfall 9 and bug M2 of the phase 2 audit in one file: two pages with a
    # real text layer, three scanned annex pages behind them. Three fifths of the
    # pages are empty, which is under the measured share, so the readable half is
    # indexed instead of the whole agreement being sent through OCR.
    outcome = extract_pdf(str(CORPUS / "14-pacht-mit-anhang.pdf"))

    assert outcome.state is State.INDEXED
    assert outcome.truncated is False
    assert "Pachtvereinbarung" in outcome.text
    assert "Der jährliche Zins" in outcome.text


def test_a_pdf_with_a_real_text_layer_is_never_treated_as_a_scan() -> None:
    # The expensive half of the asymmetry. A text PDF that lands in the OCR track
    # costs minutes of CPU per document from this phase on, so the two smallest
    # genuine text layers of the corpus are pinned here.
    for name, needle in (("01-text-layer.pdf", "real text layer"), ("09-bescheid.pdf", "Grundstücks")):
        outcome = extract_pdf(str(CORPUS / name))

        assert outcome.state is State.INDEXED, f"{name} came back as {outcome.state}/{outcome.reason}"
        assert needle in outcome.text


def test_a_cover_page_in_front_of_nine_scans_is_an_ocr_candidate_and_not_a_text_document(tmp_path: Path) -> None:
    # The case bug M2 names word for word. The document average would have to
    # decide between 200 characters and ten pages; the per page count does not
    # have to decide anything, it counts nine empty pages out of ten and hands
    # the file over. The cover text is not lost: OCR reads that page as well.
    path = _write(tmp_path, "deckblatt.pdf", _pdf_with_page_lines([_lines_for(200), *([()] * 9)]))

    outcome = extract_pdf(path)

    assert outcome == ExtractionOutcome.skipped(Reason.NO_TEXT_LAYER)


def test_two_dense_pages_in_front_of_eight_scans_are_not_ticked_off_as_a_text_document(tmp_path: Path) -> None:
    # The one the old document wide rule got wrong: 900 characters over ten pages
    # are far above 25 per page on average, so the file was indexed and the eight
    # scanned pages were never looked at again.
    path = _write(tmp_path, "zwei-und-acht.pdf", _pdf_with_page_lines([_lines_for(450), _lines_for(450), *([()] * 8)]))

    outcome = extract_pdf(path)

    assert outcome == ExtractionOutcome.skipped(Reason.NO_TEXT_LAYER)


@pytest.mark.parametrize(
    ("scanned", "readable", "expected"),
    [
        (2, 1, State.INDEXED),
        (3, 1, State.SKIPPED),
        (6, 3, State.INDEXED),
        (7, 3, State.SKIPPED),
    ],
)
def test_the_share_of_empty_pages_decides_at_the_measured_two_thirds(
    scanned: int, readable: int, expected: State, tmp_path: Path
) -> None:
    # Exactly two thirds is still a document with a text layer, above it the file
    # is a scan. The boundary is stated here rather than left to a comparison
    # operator nobody reads twice.
    pages = [*([_lines_for(200)] * readable), *([()] * scanned)]
    path = _write(tmp_path, f"anteil-{scanned}-{readable}.pdf", _pdf_with_page_lines(pages))

    outcome = extract_pdf(path)

    assert outcome.state is expected
    assert scanned / (scanned + readable) > _SCAN_PAGE_SHARE if expected is State.SKIPPED else True


@pytest.mark.parametrize(
    ("characters", "expected"),
    [(_MIN_CHARS_PER_PAGE - 1, State.SKIPPED), (_MIN_CHARS_PER_PAGE + 1, State.INDEXED)],
)
def test_one_page_counts_as_scanned_below_the_measured_threshold(
    characters: int, expected: State, tmp_path: Path
) -> None:
    # The measured neighbours of the threshold are a headline only page at 12
    # characters and the sparsest genuine text page of the corpus at 29. The
    # constant sits between them, and this pins which side of it means what.
    path = _write(tmp_path, f"seite-{characters}.pdf", _pdf_with_page_lines([(_filler(characters),)]))

    outcome = extract_pdf(path)

    assert outcome.state is expected


CORPUS_PDF_VERDICTS: dict[str, tuple[State, Reason | None]] = {
    "01-text-layer.pdf": (State.INDEXED, None),
    "02-scan-no-text-layer.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "06-zero-bytes.pdf": (State.FAILED, Reason.EMPTY_FILE),
    "07-password-protected.pdf": (State.SKIPPED, Reason.ENCRYPTED),
    "09-bescheid.pdf": (State.INDEXED, None),
    "13-ratsvorlage-scan.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "14-pacht-mit-anhang.pdf": (State.INDEXED, None),
    "15-schweiz-baubewilligung.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "16-oesterreich-mitteilung.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "24-abgeschnittener-trailer.pdf": (State.FAILED, Reason.CORRUPT),
    "25-kaputte-xref.pdf": (State.INDEXED, None),
    "26-riesige-seitenzahl.pdf": (State.FAILED, Reason.CORRUPT),
    "27-nullbytes-im-kopf.pdf": (State.FAILED, Reason.CORRUPT),
    "28-ohne-seiten.pdf": (State.FAILED, Reason.CORRUPT),
    "29-doppelt-komprimiert.pdf": (State.INDEXED, None),
    "30-nur-ein-bild.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "31-riesenformat.pdf": (State.SKIPPED, Reason.NO_TEXT_LAYER),
    "32-startxref-ins-leere.pdf": (State.INDEXED, None),
    "33-seitenbaum-zyklus.pdf": (State.FAILED, Reason.CORRUPT),
}


def test_every_pdf_of_the_corpus_keeps_the_verdict_that_is_documented_for_it() -> None:
    # testdata/CORPUS.md prints this table for a reader, and the OCR acceptance
    # of the following plans counts on it. Two of the entries are the ones that
    # could hang a run instead of ending it: 26 declares one hundred thousand
    # pages in 627 bytes and 33 is a page tree that contains itself. The time
    # bound is generous on purpose; it is there to fail rather than to hang.
    names = sorted(path.name for path in CORPUS.glob("*.pdf"))
    assert names == sorted(CORPUS_PDF_VERDICTS), "a PDF was added to the corpus without a documented verdict"

    started = time.monotonic()
    for name, (state, reason) in CORPUS_PDF_VERDICTS.items():
        outcome = extract_pdf(str(CORPUS / name))

        assert outcome.state is state, f"{name} came back as {outcome.state}/{outcome.reason}"
        assert outcome.reason is reason, f"{name} came back with reason {outcome.reason}"

    assert time.monotonic() - started < 30


def test_the_dispatcher_reaches_the_pdf_route() -> None:
    path = CORPUS / "01-text-layer.pdf"

    outcome = extract(str(path), "application/pdf", path.stat().st_size)

    assert outcome.state is State.INDEXED
    assert "Findling reference corpus" in outcome.text


# ---------------------------------------------------------------------------
# OOXML: DOCX, PPTX and XLSX
# ---------------------------------------------------------------------------

NOT_A_ZIP = b"Dies ist kein Archiv, sondern schlichter Text mit falscher Endung."


def _docx_with_a_table(directory: Path) -> str:
    """A document whose text sits in a table, which is the case python-docx hides.

    Paragraphs are the obvious half of a DOCX. A table cell is a paragraph inside
    a cell inside a row, and an extractor that only walks document.paragraphs
    loses every invoice, every list of names and every form.
    """
    document = docx.Document()
    document.add_paragraph("Aktenvermerk der Gemeinde.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Aktenzeichen"
    table.rows[0].cells[1].text = "Grundstuecksverkehr 2026-0815"
    table.rows[1].cells[0].text = "Frist"
    table.rows[1].cells[1].text = "Kuendigungsfrist drei Monate"
    target = directory / "vermerk.docx"
    document.save(str(target))
    return str(target)


def _pptx_with_two_text_frames(directory: Path) -> str:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(3657600), Emu(914400))
    first.text_frame.text = "Sitzungsvorlage des Ausschusses"
    second = slide.shapes.add_textbox(Emu(914400), Emu(2743200), Emu(3657600), Emu(914400))
    second.text_frame.text = "Beschlussvorschlag zur Strassenbaubeitragssatzung"
    target = directory / "vorlage.pptx"
    presentation.save(str(target))
    return str(target)


def _xlsx(directory: Path, rows: int, columns: int) -> str:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    for row in range(rows):
        sheet.append([f"Zelle {row}-{column}" for column in range(columns)])
    target = directory / "mappe.xlsx"
    workbook.save(str(target))
    workbook.close()
    return str(target)


def test_a_docx_yields_its_paragraphs_including_the_german_special_characters() -> None:
    outcome = extract_docx(str(CORPUS / "03-document.docx"))

    assert outcome.state is State.INDEXED
    assert "office document part" in outcome.text
    assert "Umlaute im Text: Grundstück, Ausschuss, Maßnahme." in outcome.text


def test_the_table_cells_of_a_docx_are_indexed_as_well(tmp_path: Path) -> None:
    outcome = extract_docx(_docx_with_a_table(tmp_path))

    assert outcome.state is State.INDEXED
    assert "Aktenvermerk der Gemeinde." in outcome.text
    assert "Aktenzeichen" in outcome.text
    assert "Kuendigungsfrist drei Monate" in outcome.text


def test_a_pptx_yields_the_text_of_every_shape_that_has_a_text_frame(tmp_path: Path) -> None:
    outcome = extract_pptx(_pptx_with_two_text_frames(tmp_path))

    assert outcome.state is State.INDEXED
    assert "Sitzungsvorlage des Ausschusses" in outcome.text
    assert "Beschlussvorschlag zur Strassenbaubeitragssatzung" in outcome.text


def test_a_spreadsheet_is_opened_read_only_and_data_only_and_never_otherwise(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    # Without read_only openpyxl builds the whole workbook in memory, and a single
    # export file tips a 4 GB container over. Without data_only every formula cell
    # arrives as its formula, so the index fills with =SUM(A1:A9) instead of text.
    seen: list[dict[str, object]] = []
    original = openpyxl.load_workbook

    def note(*args: object, **kwargs: object) -> object:
        seen.append(dict(kwargs))
        return original(*args, **kwargs)  # pyright: ignore[reportArgumentType]

    monkeypatch.setattr(openpyxl, "load_workbook", note)

    outcome = extract_xlsx(_xlsx(tmp_path, rows=3, columns=2))

    assert outcome.state is State.INDEXED
    assert "Zelle 2-1" in outcome.text
    assert len(seen) == 1
    assert seen[0]["read_only"] is True
    assert seen[0]["data_only"] is True


def test_a_spreadsheet_over_the_cell_limit_is_skipped_instead_of_half_read(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    # The limit is lowered instead of writing a workbook with 200000 cells: the
    # behaviour under test is the counter and the verdict skipped(too_many_cells),
    # not openpyxl's ability to write a large file.
    monkeypatch.setenv("FINDLING_MAX_CELLS", "10")
    config.settings.cache_clear()
    try:
        outcome = extract_xlsx(_xlsx(tmp_path, rows=5, columns=4))

        assert outcome == ExtractionOutcome.skipped(Reason.TOO_MANY_CELLS)
    finally:
        config.settings.cache_clear()


@pytest.mark.parametrize(
    ("extractor", "name"),
    [
        (extract_docx, "brief.docx"),
        (extract_pptx, "vorlage.pptx"),
        (extract_xlsx, "mappe.xlsx"),
    ],
)
def test_an_ooxml_name_on_a_file_that_is_no_archive_is_failed_corrupt(
    extractor: Callable[[str], ExtractionOutcome], name: str, tmp_path: Path
) -> None:
    outcome = extractor(_write(tmp_path, name, NOT_A_ZIP))

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


def test_a_truncated_docx_is_failed_corrupt_instead_of_an_escaping_exception(tmp_path: Path) -> None:
    # The measured case: a package that stops halfway raises PackageNotFoundError,
    # which reads like "there is no document here" and means "this one is broken".
    intact = (CORPUS / "03-document.docx").read_bytes()

    outcome = extract_docx(_write(tmp_path, "halb.docx", intact[: len(intact) // 2]))

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


def test_the_dispatcher_reaches_the_three_ooxml_routes(tmp_path: Path) -> None:
    documents = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": (
            _docx_with_a_table(tmp_path),
            "Aktenzeichen",
        ),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": (
            _pptx_with_two_text_frames(tmp_path),
            "Sitzungsvorlage des Ausschusses",
        ),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": (
            _xlsx(tmp_path, rows=3, columns=2),
            "Zelle 1-1",
        ),
    }

    for mime, (path, needle) in documents.items():
        outcome = extract(path, mime, Path(path).stat().st_size)

        assert outcome.state is State.INDEXED
        assert needle in outcome.text


# ---------------------------------------------------------------------------
# OpenDocument: ODT, ODS and ODP
# ---------------------------------------------------------------------------

ODF_HEADER = (
    '<?xml version="1.0" encoding="UTF-8"?>'
    '<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
    ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
    "<office:body><office:text>"
)
ODF_FOOTER = "</office:text></office:body></office:document-content>"

ODF_BODY = (
    "<text:h>Sitzungsvorlage des Ausschusses</text:h>"
    "<text:p>Erster Absatz mit einer <text:span>Anlage</text:span>.</text:p>"
    "<text:p>Zweiter Absatz zur Grundstücksverkehrsgenehmigung.</text:p>"
)


def _odf(directory: Path, name: str, content: str | None) -> str:
    """An OpenDocument file, written with zipfile so its innards stay readable.

    Building the fixture here rather than adding three binaries to the reference
    corpus keeps that corpus at the size it was designed for and puts the whole
    document into the diff of this test.
    """
    target = directory / name
    with ZipFile(target, "w") as archive:
        archive.writestr("mimetype", "application/vnd.oasis.opendocument.text")
        if content is not None:
            archive.writestr("content.xml", content)
    return str(target)


def test_an_odt_yields_every_heading_and_paragraph_in_document_order(tmp_path: Path) -> None:
    outcome = extract_odf(_odf(tmp_path, "vorlage.odt", ODF_HEADER + ODF_BODY + ODF_FOOTER))

    assert outcome.state is State.INDEXED
    assert outcome.text.index("Sitzungsvorlage") < outcome.text.index("Erster Absatz")
    assert outcome.text.index("Erster Absatz") < outcome.text.index("Zweiter Absatz")
    assert "Grundstücksverkehrsgenehmigung" in outcome.text


def test_the_parts_are_joined_with_a_space_instead_of_being_glued_together(tmp_path: Path) -> None:
    # Two paragraphs concatenated produce a word that exists in neither of them,
    # and that word is then the only thing the index knows about the boundary.
    outcome = extract_odf(_odf(tmp_path, "vorlage.odt", ODF_HEADER + ODF_BODY + ODF_FOOTER))

    assert "Anlage. Zweiter Absatz" in outcome.text
    assert "Anlage.Zweiter" not in outcome.text


def test_a_declared_decompression_bomb_is_skipped_before_it_is_read(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    # The archive directory names the real size; reading first would BE the
    # attack (security audit M4: 815 kB ODT -> 800 MB content.xml). The cap is
    # lowered so the fixture stays a fixture instead of a real bomb.
    monkeypatch.setattr("findling.extract.odf.EXTRACT_ARCHIVE_MEMBER_MAX_BYTES", 64)

    outcome = extract_odf(_odf(tmp_path, "bombe.odt", "x" * 65))

    assert outcome.state is State.SKIPPED
    assert outcome.reason is Reason.TOO_LARGE


def test_an_office_package_with_an_oversized_part_is_skipped_before_any_loader_runs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr("findling.extract.office.EXTRACT_ARCHIVE_MEMBER_MAX_BYTES", 64)
    target = tmp_path / "bombe.docx"
    with ZipFile(target, "w") as archive:
        archive.writestr("word/document.xml", "x" * 65)

    # The archive is not even a valid package; a verdict of too_large therefore
    # proves the guard ran before the loader ever touched it.
    for extractor in (extract_docx, extract_pptx, extract_xlsx):
        outcome = extractor(str(target))

        assert outcome.state is State.SKIPPED
        assert outcome.reason is Reason.TOO_LARGE


@pytest.mark.parametrize("name", ["tabelle.ods", "praesentation.odp"])
def test_ods_and_odp_run_through_the_very_same_path(name: str, tmp_path: Path) -> None:
    outcome = extract_odf(_odf(tmp_path, name, ODF_HEADER + ODF_BODY + ODF_FOOTER))

    assert outcome.state is State.INDEXED
    assert "Sitzungsvorlage des Ausschusses" in outcome.text


def test_an_archive_without_a_content_part_is_failed_corrupt(tmp_path: Path) -> None:
    outcome = extract_odf(_odf(tmp_path, "leer.odt", None))

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


def test_a_file_that_is_no_archive_at_all_is_failed_corrupt(tmp_path: Path) -> None:
    outcome = extract_odf(_write(tmp_path, "kein-archiv.odt", NOT_A_ZIP))

    assert outcome == ExtractionOutcome.failed(Reason.CORRUPT)


def test_a_broken_content_part_is_failed_xml_invalid_and_not_corrupt(tmp_path: Path) -> None:
    # The two reasons are different information for whoever reads the status page:
    # a package that cannot be opened is a different repair job from a package
    # that opens and holds nonsense.
    outcome = extract_odf(_odf(tmp_path, "kaputt.odt", ODF_HEADER + "<text:p>Ohne Ende"))

    assert outcome == ExtractionOutcome.failed(Reason.XML_INVALID)


def test_the_content_parser_follows_neither_an_entity_nor_a_url(tmp_path: Path) -> None:
    # An OpenDocument file is untrusted input, and its content part is XML. Without
    # the three switches a crafted document makes this container read a local file
    # and call an address of the attacker's choosing, and both results then travel
    # into the search index.
    secret = tmp_path / "passphrase.txt"
    secret.write_bytes(b"TOPSECRET-PASSPHRASE")
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<!DOCTYPE office:document-content [\n"
        f'  <!ENTITY leak SYSTEM "{secret.as_uri()}">\n'
        '  <!ENTITY call SYSTEM "http://127.0.0.1:9/collect">\n'
        "]>"
        '<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        "<office:body><office:text><text:p>Anlage &leak; &call;</text:p></office:text></office:body>"
        "</office:document-content>"
    )
    path = _odf(tmp_path, "xxe.odt", content)

    outcome = extract_odf(path)

    assert "TOPSECRET" not in outcome.text
    assert "PASSPHRASE" not in outcome.text

    # The counter measurement, so the test proves the switches rather than the
    # fixture: the same file entity through a parser built without the switches
    # does leak. The network entity stays out of this parse on purpose: with
    # no_network=True lxml on Linux raises XMLSyntaxError at the load attempt
    # instead of parsing on (Windows parses on), and this measurement is about
    # the leak, not about the refusal.
    file_only_content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<!DOCTYPE office:document-content [\n"
        f'  <!ENTITY leak SYSTEM "{secret.as_uri()}">\n'
        "]>"
        '<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
        ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        "<office:body><office:text><text:p>Anlage &leak;</text:p></office:text></office:body>"
        "</office:document-content>"
    )
    greedy = etree.XMLParser(resolve_entities=True, no_network=True, load_dtd=True)
    leaked = etree.fromstring(file_only_content.encode("utf-8"), parser=greedy)

    assert "TOPSECRET-PASSPHRASE" in "".join(leaked.itertext())


def test_the_archive_is_only_ever_read_and_never_unpacked(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    # extractall() writes attacker chosen names to the file system, which is how
    # zip slip works. The one part that is needed is read into memory instead, and
    # this test is the reason that cannot be quietly changed back.
    def unreachable(*_args: object, **_kwargs: object) -> None:
        raise AssertionError("the archive was unpacked to disk")

    monkeypatch.setattr(ZipFile, "extractall", unreachable)

    outcome = extract_odf(_odf(tmp_path, "vorlage.odt", ODF_HEADER + ODF_BODY + ODF_FOOTER))

    assert outcome.state is State.INDEXED


def test_every_route_of_the_allowlist_has_an_extractor_behind_it(tmp_path: Path) -> None:
    # The dispatcher had five open places while the document formats were missing.
    # This walks the whole allowlist, so a route that loses its extractor cannot
    # hide behind the fourteen mimetypes that still work.
    odf_document = _odf(tmp_path, "vorlage.odt", ODF_HEADER + ODF_BODY + ODF_FOOTER)
    fixtures = {
        Route.PDF: str(CORPUS / "01-text-layer.pdf"),
        Route.DOCX: _docx_with_a_table(tmp_path),
        Route.PPTX: _pptx_with_two_text_frames(tmp_path),
        Route.XLSX: _xlsx(tmp_path, rows=3, columns=2),
        Route.ODF: odf_document,
        Route.HTML: _write(tmp_path, "seite.html", b"<html><body><p>Ein Absatz.</p></body></html>"),
        Route.RTF: _write(tmp_path, "brief.rtf", rb"{\rtf1\ansi Ein Absatz.\par}"),
        Route.PLAIN: _write(tmp_path, "notiz.txt", b"Ein Absatz."),
    }

    # Route.OCR is the one route with no fixture here, and that is the assertion
    # rather than an omission: no mimetype maps to it, it is reached by the kind
    # of the job (plan 03-09), and this loop walks the allowlist. Its extractor is
    # covered in tests/test_ocr.py, where the engine can be stood in for.
    assert set(fixtures) == set(Route) - {Route.OCR}
    assert Route.OCR not in set(ALLOWED_MIMETYPES.values())

    for mime, route in ALLOWED_MIMETYPES.items():
        path = fixtures[route]

        outcome = extract(path, mime, Path(path).stat().st_size)

        assert outcome.state is State.INDEXED, f"{mime} came back as {outcome.state}/{outcome.reason}"

"""DOCX, PPTX and XLSX: three ZIP packages, and one of them is dangerous.

The two word processing formats are unremarkable. The spreadsheet is not, and
this module exists in the shape it has because of it.

**The workbook is opened read only and with values only, never otherwise.**
Without the read only switch openpyxl builds the entire workbook as objects in
memory before a single cell is read, and one export file with a million rows is
the classic way to tip a container with 4 GB of RAM over. Without the values only
switch a formula cell arrives as its formula, so the index fills up with
``=SUM(A1:A9)`` instead of the number a person would search for. On top of both
sits a hard cell limit that is counted while reading and ends the loop:
skipped(too_many_cells) is a verdict a user can act on, half a spreadsheet in the
index is not.

**Known gap, deliberately left open.** python-docx reads the document body, and
that leaves out headers, footers, footnotes and text boxes. Reading them back
would mean parsing ``word/header*.xml`` and ``word/footnotes.xml`` with lxml,
which is a piece of work with its own error paths and is not part of phase 2.
For a search index the body is where the words are; documented non support is
more honest than a half finished second parser. Tables inside tables are the
same kind of border: the cells of a nested table are not walked.

**No taxonomy of its own.** A broken package raises, and the exception is handed
to ExtractionOutcome.from_exception, which is the single place that knows
exception classes. A format module that invented its own reason codes would be
the first step towards two lists that disagree.
"""

from __future__ import annotations

from zipfile import BadZipFile, ZipFile

import docx
import openpyxl
import pptx
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFound
from pptx.exc import PackageNotFoundError as PptxPackageNotFound

from findling import config
from findling.config import EXTRACT_ARCHIVE_MEMBER_MAX_BYTES
from findling.extract.dispatch import cap_text
from findling.extract.errors import ExtractionOutcome, Reason

# What a package that is not a package raises, measured against a file that is no
# archive at all and against a truncated one. Both go to from_exception rather
# than to a reason picked here.
_BROKEN_PACKAGE = (BadZipFile, DocxPackageNotFound, PptxPackageNotFound)


def _oversized_part(path: str) -> ExtractionOutcome | None:
    """The bomb check, before any loader touches the package.

    The archive directory declares the uncompressed size of every part, and the
    loaders below read whole parts into memory; a bomb therefore has to be
    refused on the declaration, because after the read it is the attack
    (security audit M4). zipfile enforces the declared size on read, so the
    declaration cannot be lied past. A package that will not even open is left
    to the loader, whose exception carries the better diagnosis.
    """
    try:
        with ZipFile(path) as archive:
            if any(info.file_size > EXTRACT_ARCHIVE_MEMBER_MAX_BYTES for info in archive.infolist()):
                return ExtractionOutcome.skipped(Reason.TOO_LARGE)
    except (BadZipFile, OSError):
        return None
    return None


def extract_docx(path: str) -> ExtractionOutcome:
    """The body text of a Word document: paragraphs first, then table cells.

    Defined at module level, like every extractor here, so it survives the
    process boundary of the extraction child.
    """
    oversized = _oversized_part(path)
    if oversized is not None:
        return oversized

    try:
        document = docx.Document(path)
    except _BROKEN_PACKAGE as error:
        return ExtractionOutcome.from_exception(error)

    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts += [cell.text for cell in row.cells]
    return cap_text("\n".join(part for part in parts if part.strip()))


def extract_pptx(path: str) -> ExtractionOutcome:
    """Every shape of every slide that carries a text frame.

    Shapes without a text frame are pictures, lines and grouping boxes. Asking
    them for text raises rather than returning nothing, which is why the question
    is asked before the text is taken.
    """
    oversized = _oversized_part(path)
    if oversized is not None:
        return oversized

    try:
        presentation = pptx.Presentation(path)
    except _BROKEN_PACKAGE as error:
        return ExtractionOutcome.from_exception(error)

    parts: list[str] = []
    for slide in presentation.slides:
        # has_text_frame is the guard python-pptx offers, and it is a runtime one:
        # the base class of a shape does not declare text_frame, so the type
        # checker cannot see that the attribute exists once the flag is true. The
        # ignore is as narrow as the problem, and the alternative, an isinstance
        # test against the concrete shape classes, would silently drop the
        # placeholder types that carry most of the text on a slide.
        parts += [
            shape.text_frame.text  # pyright: ignore[reportAttributeAccessIssue]
            for shape in slide.shapes
            if shape.has_text_frame
        ]
    return cap_text("\n".join(part for part in parts if part.strip()))


def extract_xlsx(path: str) -> ExtractionOutcome:
    """Cell values of every sheet, bounded by the cell limit.

    The two keyword arguments are not tuning, they are the reason this function
    can run in a small container at all; the module docstring above says what
    happens without them. The workbook is closed in a finally, because the read
    only mode keeps file handles on the archive open.
    """
    oversized = _oversized_part(path)
    if oversized is not None:
        return oversized

    limit = config.settings().max_cells
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except _BROKEN_PACKAGE as error:
        return ExtractionOutcome.from_exception(error)

    parts: list[str] = []
    seen = 0
    try:
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                seen += len(row)
                if seen > limit:
                    # Stop here rather than deliver half a workbook. A truncated
                    # spreadsheet in the index looks like a complete one.
                    return ExtractionOutcome.skipped(Reason.TOO_MANY_CELLS)
                parts += [str(value) for value in row if value is not None]
    finally:
        workbook.close()

    return cap_text("\n".join(part for part in parts if part.strip()))
